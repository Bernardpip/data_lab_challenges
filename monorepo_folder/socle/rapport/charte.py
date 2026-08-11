"""Charte du rapport PowerPoint — dérivée des tokens de l'écran.

Le pilote redéclarait ses couleurs ici, en commentant « reprise de
components/tokens.py ». Une copie ne se met pas à jour : changer la teinte du
slot 1 dans les tokens laissait le PPTX sur l'ancienne, et les deux livrables
cessaient silencieusement de se ressembler. Les valeurs sont donc CONVERTIES
depuis `socle.design.tokens`, jamais retapées.

Les graphiques sont NATIFS PowerPoint (`add_chart`), donc éditables et
redimensionnables par le lecteur — pas des images figées.
"""

# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
# pyrefly: ignore [missing-import]
from pptx.enum.text import PP_ALIGN
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt

from socle.design.tokens import COLORS, SERIES, STATUS


def rgb(hexa):
    """« #2a78d6 » → RGBColor. Le point d'entrée de toute couleur du rapport.

    Une couleur de marque propre au défi (le vert d'un drapeau, par exemple)
    se déclare avec ce convertisseur dans le fichier du défi, pas ici : le
    socle ne connaît aucun commanditaire.
    """

    valeur = hexa.lstrip("#")

    return RGBColor(int(valeur[0:2], 16), int(valeur[2:4], 16), int(valeur[4:6], 16))


ENCRE = rgb(COLORS["text"])
ENCRE_SECONDAIRE = rgb(COLORS["textSecondary"])
MUTED = rgb(COLORS["textMuted"])
PRIMAIRE = rgb(COLORS["primary"])
SURFACE = rgb(COLORS["surface"])
ROUGE = rgb(STATUS["critical"])

# Mêmes slots que l'écran, dans le même ordre : une série bleue dans le
# tableau de bord doit rester bleue dans le rapport.
SERIE = [rgb(couleur) for couleur in SERIES]
SERIE_1, SERIE_2 = SERIE[0], SERIE[1]

# 16:9, format par défaut des vidéoprojecteurs.
LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

POLICE = "Calibri"


def texte(cadre, contenu, taille=14, gras=False, couleur=ENCRE,
          espace_apres=6, aligne=PP_ALIGN.LEFT, premier=False):
    """Un paragraphe stylé dans un cadre de texte.

    `premier=True` réutilise le paragraphe que python-pptx crée d'office —
    sans quoi chaque cadre commencerait par une ligne vide.
    """

    paragraphe = cadre.paragraphs[0] if premier else cadre.add_paragraph()
    paragraphe.alignment = aligne
    paragraphe.space_after = Pt(espace_apres)

    run = paragraphe.add_run()
    run.text = contenu
    run.font.size = Pt(taille)
    run.font.bold = gras
    run.font.color.rgb = couleur
    run.font.name = POLICE

    return paragraphe


def titre_page(slide, numero, titre, sous_titre=None):
    """Bandeau de titre commun à toutes les pages."""

    bande = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    cadre = bande.text_frame
    cadre.word_wrap = True

    texte(cadre, f"{numero} · {titre}", taille=26, gras=True, premier=True,
          espace_apres=2)

    if sous_titre:
        texte(cadre, sous_titre, taille=13, couleur=MUTED, espace_apres=0)

    # Filet indigo sous le titre.
    ligne = slide.shapes.add_textbox(Inches(0.6), Inches(1.32), Inches(12.1), Inches(0.03))
    ligne.fill.solid()
    ligne.fill.fore_color.rgb = PRIMAIRE
    ligne.line.fill.background()


def pied(slide, numero_page, lg):
    """Pied de page — le libellé vient du domaine i18n du défi."""

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.35))
    cadre = boite.text_frame
    paragraphe = cadre.paragraphs[0]

    run = paragraphe.add_run()
    run.text = lg.t("pied", {"page": numero_page})
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED
    run.font.name = POLICE


def bloc_constat(slide, gauche, haut, largeur, hauteur, valeur, libelle,
                 detail, couleur=PRIMAIRE):
    """Tuile de chiffre-clé, alignée sur les `stat_tiles` du tableau de bord."""

    boite = slide.shapes.add_textbox(gauche, haut, largeur, hauteur)
    cadre = boite.text_frame
    cadre.word_wrap = True

    texte(cadre, libelle, taille=11, couleur=MUTED, espace_apres=2, premier=True)
    texte(cadre, valeur, taille=32, gras=True, couleur=couleur, espace_apres=2)
    texte(cadre, detail, taille=10, couleur=ENCRE_SECONDAIRE, espace_apres=0)


def bloc_analyse(slide, gauche, haut, largeur, hauteur, nom, question,
                 resultat, lecture, couleur_resultat=ENCRE):
    """Bloc d'analyse : la question posée, le résultat obtenu, sa portée.

    Volontairement SANS notation mathématique : le livrable attendu explique
    les analyses et leurs résultats, pas l'appareil de calcul — celui-ci reste
    consultable dans le tableau de bord.
    """

    boite = slide.shapes.add_textbox(gauche, haut, largeur, hauteur)
    cadre = boite.text_frame
    cadre.word_wrap = True

    texte(cadre, nom, taille=12, gras=True, couleur=PRIMAIRE, espace_apres=2,
          premier=True)
    texte(cadre, question, taille=10, couleur=MUTED, espace_apres=5)
    texte(cadre, resultat, taille=13, gras=True, couleur=couleur_resultat,
          espace_apres=4)
    texte(cadre, lecture, taille=9.5, couleur=ENCRE_SECONDAIRE, espace_apres=0)


def style_graphe(graphique, titre=None):
    """Typographie d'un graphe natif, accordée au reste du document."""

    graphique.has_title = titre is not None

    if titre:
        graphique.chart_title.text_frame.text = titre
        graphique.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    graphique.font.size = Pt(10)
    graphique.font.color.rgb = ENCRE_SECONDAIRE
    graphique.font.name = POLICE
