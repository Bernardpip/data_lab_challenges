"""Génère le rapport PowerPoint (10 pages maximum) attendu comme livrable.

    python scripts/generer_presentation.py            # les deux langues
    python scripts/generer_presentation.py fr         # une seule

Ce fichier ne porte que **les 10 pages** — ce que ce défi-ci raconte. Le
bandeau de titre, le pied, la tuile de chiffre-clé, le bloc d'analyse, le
style des graphes natifs, l'objet `Langue` et l'assemblage bilingue viennent
de `socle.rapport`.

Trois choix qui ne se renégocient pas :

  · les graphiques sont NATIFS PowerPoint (`add_chart`), donc éditables et
    redimensionnables par le lecteur — pas des images figées ;
  · tous les chiffres sont recalculés par `collecter()` depuis les MÊMES
    fonctions que le tableau de bord : la présentation ne peut pas diverger
    de l'écran ;
  · `collecter()` est appelé UNE seule fois pour les deux langues. Deux
    chargements laisseraient les versions diverger sur un arrondi.

Les 10 pages de la trame :

    1. Couverture                    6. Économétrie
    2. Données, nettoyage, démarche  7. Moyens publics
    3. Le constat structurant        8. Le croisement attendu (et ses limites)
    4. Sa dynamique dans le temps    9. Leviers d'action
    5. L'effet mis en évidence      10. Conclusion et limites
"""

import sys
from pathlib import Path

RACINE = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(RACINE))

# pyrefly: ignore [missing-import]
from pptx.chart.data import CategoryChartData                   # noqa: E402
# pyrefly: ignore [missing-import]
from pptx.enum.chart import XL_CHART_TYPE                       # noqa: E402
# pyrefly: ignore [missing-import]
from pptx.util import Inches                                    # noqa: E402

from socle import i18n                                          # noqa: E402

i18n.configurer(RACINE / "i18n" / "locales")

from socle.i18n import LANGUES                                  # noqa: E402
from socle.rapport import charte, generer_toutes                # noqa: E402



def page_1_couverture(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    boite = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5),
                                     Inches(3.0))
    cadre = boite.text_frame
    cadre.word_wrap = True

    charte.texte(cadre, lg.t("p1_titre"), taille=40, gras=True, premier=True,
                 espace_apres=10)
    charte.texte(cadre, lg.t("p1_sous_titre"), taille=16,
                 couleur=charte.ENCRE_SECONDAIRE, espace_apres=22)
    charte.texte(cadre, lg.t("p1_accroche", {
        "cantons": lg.nb(c["cantons"]),
        "part": lg.nb(c["part_pop_exposee"], 0),
        "distance": lg.nb(c["distance_mediane"], 0),
    }), taille=13, couleur=charte.ENCRE, espace_apres=26)
    charte.texte(cadre, lg.t("p1_auteur"), taille=12, couleur=charte.MUTED)


def page_2_demarche(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 2, lg.t("p2_titre"), lg.t("p2_sous_titre"))

    tuiles = [
        (lg.nb(c["jeux"]), "p2_t1", "p2_t1_d"),
        (lg.nb(c["cantons"]), "p2_t2", "p2_t2_d"),
        (lg.nb(c["ouvrages"]), "p2_t3", "p2_t3_d"),
        (f'{c["publies"]} / {c["decrits"]}', "p2_t4", "p2_t4_d"),
    ]

    for index, (valeur, libelle, detail) in enumerate(tuiles):
        charte.bloc_constat(
            slide, Inches(0.6 + index * 3.1), Inches(1.85), Inches(2.9),
            Inches(1.3), valeur=valeur, libelle=lg.t(libelle),
            detail=lg.t(detail),
        )

    for index, cle in enumerate(("p2_etape_1", "p2_etape_2", "p2_etape_3")):
        charte.bloc_analyse(
            slide, Inches(0.6 + index * 4.15), Inches(3.5), Inches(3.9),
            Inches(2.6), nom=lg.t(f"{cle}_nom"), question=lg.t(f"{cle}_q"),
            resultat=lg.t(f"{cle}_r"), lecture=lg.t(f"{cle}_l"),
        )

    charte.pied(slide, 2, lg)


def page_3_distance(prs, c, lg):
    """Objectif 1 — cartographier la répartition : à quelle distance est l'eau.

    La première question du sujet n'est pas « combien d'ouvrages ? » mais « où
    sont-ils ? ». La réponse tient à une distance, et elle est mesurée sur la
    projection métrique du pays, ouvrage par ouvrage.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 3, lg.t("p3_titre"), lg.t("p3_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=[lg.t("p3_r1"), lg.t("p3_r2"), lg.t("p3_r3")],
            valeurs=c["rayons_part"], titre=lg.t("p3_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(2.0), Inches(4.3), Inches(1.4),
        valeur=f'{lg.nb(c["distance_mediane"], 0)} km', libelle=lg.t("p3_t1"),
        detail=lg.t("p3_t1_d"), couleur=charte.DANGER,
    )
    charte.bloc_constat(
        slide, Inches(8.3), Inches(3.5), Inches(4.3), Inches(1.4),
        valeur=lg.nb(c["cantons_loin"]), libelle=lg.t("p3_t2"),
        detail=lg.t("p3_t2_d", {"part": lg.nb(c["part_pop_loin"], 0)}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(5.0), Inches(4.3), Inches(1.5),
        nom=lg.t("p3_lecture_nom"), question=lg.t("p3_lecture_q"),
        resultat=lg.t("p3_lecture_r", {"r": lg.nb(c["clark_evans"], 2)}),
        lecture=lg.t("p3_lecture_l"),
    )

    charte.pied(slide, 3, lg)


def page_4_etat(prs, c, lg):
    """Objectif 2 — le taux de fonctionnalité, et pourquoi il ne se calcule pas.

    C'est la page la plus importante du rapport, et la seule qui répond « non ».
    Le sujet demande des taux de panne par région ; aucun champ du corpus ne
    porte l'état d'un ouvrage. Deux substituts existent, ils sont nommés comme
    tels, et ce qu'ils ne remplacent pas est écrit.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 4, lg.t("p4_titre"), lg.t("p4_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=c["maintenance_regions"], valeurs=c["maintenance_parts"],
            titre=lg.t("p4_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(2.0), Inches(4.3), Inches(1.4),
        valeur=lg.nb(c["absents"]), libelle=lg.t("p4_t1"),
        detail=lg.t("p4_t1_d", {"decrits": c["decrits"]}),
        couleur=charte.DANGER,
    )
    charte.bloc_constat(
        slide, Inches(8.3), Inches(3.5), Inches(4.3), Inches(1.4),
        valeur=lg.nb(c["en_attente"]), libelle=lg.t("p4_t2"),
        detail=lg.t("p4_t2_d", {"delai": lg.nb(c["delai_remise"], 0)}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(5.0), Inches(4.3), Inches(1.5),
        nom=lg.t("p4_lecture_nom"), question=lg.t("p4_lecture_q"),
        resultat=lg.t("p4_lecture_r"), lecture=lg.t("p4_lecture_l"),
    )

    charte.pied(slide, 4, lg)


def page_5_couverture(prs, c, lg):
    """La pression démographique face à l'inventaire."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 5, lg.t("p5_titre"), lg.t("p5_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=c["regions"], valeurs=c["regions_hab_ouvrage"],
            titre=lg.t("p5_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(2.0), Inches(4.3), Inches(1.4),
        valeur=lg.compact(c["plateaux_hab_ouvrage"]), libelle=lg.t("p5_t1"),
        detail=lg.t("p5_t1_d"), couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(3.6), Inches(4.3), Inches(2.8),
        nom=lg.t("p5_lecture_nom"), question=lg.t("p5_lecture_q"),
        resultat=lg.t("p5_lecture_r", {
            "national": lg.compact(c["national_hab_ouvrage"])}),
        lecture=lg.t("p5_lecture_l"),
    )

    charte.pied(slide, 5, lg)


def page_6_risque(prs, c, lg):
    """Objectif 4 — croiser le risque d'inondation et les ouvrages.

    Deux résultats sur la même page, parce qu'ils se répondent : le risque est
    CONCENTRÉ — dix-sept cantons, un tiers du pays —, et il n'explique pas qui
    reçoit un ouvrage. Le troisième bloc désamorce l'objection qui vient
    aussitôt : l'indice officiel ne classerait-il pas simplement les cantons
    peuplés ? Non, et c'est mesuré.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 6, lg.t("p6_titre"), lg.t("p6_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=[lg.t(f"classe_{cle}") for cle in c["classes_cles"]],
            valeurs=c["classes_population"], titre=lg.t("p6_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(1.9), Inches(4.3), Inches(1.3),
        valeur=lg.nb(c["cantons_exposes"]), libelle=lg.t("p6_t1"),
        detail=lg.t("p6_t1_d", {
            "part": lg.nb(c["part_cantons_exposes"], 0),
            "population": lg.compact(c["population_exposee"])}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(3.35), Inches(4.3), Inches(1.5),
        nom=lg.t("p6_a_nom"), question=lg.t("p6_a_q"),
        resultat=lg.t("p6_a_r", {
            "sans": lg.nb(100 * c["r2_besoin"], 0),
            "avec": lg.nb(100 * c["r2_region"], 0)}),
        lecture=lg.t("p6_a_l", {
            "cantons": lg.nb(c["prioritaires"]),
            "population": lg.compact(c["population_prioritaire"])}),
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(5.0), Inches(4.3), Inches(1.5),
        nom=lg.t("p6_b_nom"), question=lg.t("p6_b_q"),
        resultat=lg.t("p6_b_r", {"rho": lg.nb(c["rho_ampute"], 2)}),
        lecture=lg.t("p6_b_l", {
            "places": lg.nb(c["deplacement_median"], 0),
            "cantons": lg.nb(c["fri_reproductible"])}),
    )

    charte.pied(slide, 6, lg)


def page_7_cout(prs, c, lg):
    """Le résultat qui commande les autres : le coût est plat, la dotation non.

    Deux pages n'en font qu'une, et elles y gagnent : le coût par bénéficiaire
    varie de un à quinze SANS que le coût d'un ouvrage bouge, et l'argent ne
    suit pas les habitants. Le premier fait est une variable de décision ; le
    second dit qu'on ne s'en sert pas.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 7, lg.t("p7_titre"), lg.t("p7_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=[lg.t(f"p7_q{index}") for index in range(1, 6)],
            valeurs=c["cout_par_quintile"], titre=lg.t("p7_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(1.9), Inches(4.3), Inches(1.3),
        valeur=lg.compact(c["cout_median"]), libelle=lg.t("p7_t1"),
        detail=lg.t("p7_t1_d"),
    )
    charte.bloc_constat(
        slide, Inches(8.3), Inches(3.35), Inches(4.3), Inches(1.3),
        valeur=lg.nb(c["elasticite"], 2), libelle=lg.t("p7_t2"),
        detail=lg.t("p7_t2_d", {"n": c["cantons_dotes"],
                                "gini": lg.nb(c["gini"], 2)}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(4.8), Inches(4.3), Inches(1.7),
        nom=lg.t("p7_lecture_nom"), question=lg.t("p7_lecture_q"),
        resultat=lg.t("p7_lecture_r", {"rapport": lg.nb(c["ecart_cout"], 0)}),
        lecture=lg.t("p7_lecture_l", {
            "ecart": lg.nb(c["ecart_appel_offres"], 0),
            "paye": lg.compact(c["budget_paye"])}),
    )

    charte.pied(slide, 7, lg)


def page_8_rattrapage(prs, c, lg):
    """Objectif 5, premier volet — combien d'ouvrages manquent, et à quel prix.

    Trois scénarios plutôt qu'un chiffre : le nombre d'ouvrages à construire
    dépend entièrement du seuil de desserte qu'on se donne, et ce seuil est un
    CHOIX politique, pas un résultat. Le prix unitaire, lui, est observé — la
    médiane de ce que le COSO a réellement payé.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 8, lg.t("p8_titre"), lg.t("p8_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=[lg.t("p8_n1"), lg.t("p8_n2"), lg.t("p8_n3")],
            valeurs=c["scenarios_manquants"], titre=lg.t("p8_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(1.9), Inches(4.3), Inches(1.3),
        valeur=lg.nb(c["manquants_mediane"]), libelle=lg.t("p8_t1"),
        detail=lg.t("p8_t1_d", {"cout": lg.compact(c["cout_mediane"])}),
    )
    charte.bloc_constat(
        slide, Inches(8.3), Inches(3.35), Inches(4.3), Inches(1.3),
        valeur=lg.compact(c["scenarios_cout"][1]), libelle=lg.t("p8_t2"),
        detail=lg.t("p8_t2_d", {
            "ouvrages": lg.nb(c["scenarios_manquants"][1])}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(4.8), Inches(4.3), Inches(1.7),
        nom=lg.t("p8_lecture_nom"), question=lg.t("p8_lecture_q"),
        resultat=lg.t("p8_lecture_r", {
            "unitaire": lg.compact(c["cout_unitaire"])}),
        lecture=lg.t("p8_lecture_l", {"observations": c["observations_cout"]}),
    )

    charte.pied(slide, 8, lg)


def page_9_programmes(prs, c, lg):
    """Objectif 5, second volet — deux programmes datés, et leur entretien.

    Une recommandation qui ne porte ni montant ni calendrier n'engage personne.
    Ces deux-ci portent les deux, canton par canton, et disent lequel de leurs
    nombres est OBSERVÉ et lequel est POSÉ : le coût d'un forage vient du
    corpus, celui d'un aménagement contre les crues n'y figure nulle part.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 9, lg.t("p9_titre"), lg.t("p9_sous_titre"))

    _barres(slide, prs, lg, Inches(0.6), Inches(1.8), Inches(7.3), Inches(4.6),
            categories=[str(annee) for annee in c["programme_annees"]],
            valeurs=c["programme_montants"], titre=lg.t("p9_graphe"))

    charte.bloc_constat(
        slide, Inches(8.3), Inches(1.9), Inches(4.3), Inches(1.3),
        valeur=lg.compact(c["programme_total"]), libelle=lg.t("p9_t1"),
        detail=lg.t("p9_t1_d", {
            "ouvrages": lg.nb(c["programme_ouvrages"]),
            "horizon": c["programme_horizon"],
            "norme": lg.nb(c["programme_norme"])}),
    )
    charte.bloc_constat(
        slide, Inches(8.3), Inches(3.35), Inches(4.3), Inches(1.3),
        valeur=lg.compact(c["inondation_total"]), libelle=lg.t("p9_t2"),
        detail=lg.t("p9_t2_d", {
            "cantons": lg.nb(c["inondation_cantons"]),
            "horizon": c["inondation_horizon"],
            "unitaire": lg.compact(c["inondation_unitaire"])}),
        couleur=charte.DANGER,
    )
    charte.bloc_analyse(
        slide, Inches(8.3), Inches(4.8), Inches(4.3), Inches(1.7),
        nom=lg.t("p9_lecture_nom"), question=lg.t("p9_lecture_q"),
        resultat=lg.t("p9_lecture_r", {
            "part": lg.nb(c["entretien_part_observee"], 0)}),
        lecture=lg.t("p9_lecture_l", {
            "provision": lg.nb(c["entretien_taux_observe"], 2),
            "haut": lg.nb(c["entretien_part_haute"], 0)}),
    )

    charte.pied(slide, 9, lg)


def page_10_conclusion(prs, c, lg):
    """Ce que le corpus établit, et ce qu'il refuse d'établir."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, 10, lg.t("p10_titre"), lg.t("p10_sous_titre"))

    for index, cle in enumerate(("p10_a", "p10_b")):
        charte.bloc_analyse(
            slide, Inches(0.6 + index * 6.3), Inches(1.9), Inches(6.0),
            Inches(2.4), nom=lg.t(f"{cle}_nom"), question=lg.t(f"{cle}_q"),
            resultat=lg.t(f"{cle}_r"), lecture=lg.t(f"{cle}_l"),
        )

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(12.1),
                                     Inches(2.0))
    cadre = boite.text_frame
    cadre.word_wrap = True
    charte.texte(cadre, lg.t("p10_limites_titre"), taille=12, gras=True,
                 couleur=charte.PRIMAIRE, premier=True, espace_apres=6)

    for index in range(1, 4):
        charte.texte(cadre, lg.t(f"p10_limite_{index}"), taille=10,
                     couleur=charte.ENCRE_SECONDAIRE, espace_apres=4)

    charte.pied(slide, 10, lg)


def _barres(slide, prs, lg, gauche, haut, largeur, hauteur, categories,
            valeurs, titre):
    """Un graphe à barres NATIF — éditable par le lecteur, pas une image.

    Regroupé ici plutôt que répété dans cinq pages : le seul écart entre
    elles est le couple (catégories, valeurs), et une image figée priverait
    le destinataire du droit de retoucher la figure dans son propre document.
    """

    donnees = CategoryChartData()
    donnees.categories = categories
    donnees.add_series(titre, tuple(valeurs))

    cadre = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, gauche, haut, largeur, hauteur, donnees
    )
    charte.style_graphe(cadre.chart, titre)
    cadre.chart.has_legend = False

    return cadre.chart


PAGES = [
    page_1_couverture,
    page_2_demarche,
    page_3_distance,
    page_4_etat,
    page_5_couverture,
    page_6_risque,
    page_7_cout,
    page_8_rattrapage,
    page_9_programmes,
    page_10_conclusion,
]


def collecter():
    """Tous les chiffres du rapport, calculés depuis les jeux nettoyés.

    Appelée UNE fois et partagée entre les deux langues. Chaque valeur vient
    d'`utils.analytics`, d'`utils.econometrie` ou d'`utils.perimetre`, jamais
    d'un calcul écrit ici : un chiffre recalculé à deux endroits finit
    toujours par diverger, et c'est le tableau de bord qui fait foi.
    """

    from utils.data import datasets
    from utils import accessibilite, analytics, econometrie, perimetre

    data = datasets()
    cantons, tde, coso = data["cantons"], data["tde"], data["coso"]

    faits = analytics.synthese(cantons, tde, coso, data["ventes"])
    classes = analytics.population_par_classe(cantons)
    matrice = analytics.matrice_risque_equipement(cantons, tde, coso)
    hautes = classes[classes["classe_officielle"].isin(analytics.CLASSES_HAUTES)]
    prioritaires = analytics.cantons_prioritaires(cantons, tde, coso)

    cout = econometrie.fonction_de_cout(coso)
    elasticite = econometrie.elasticite_investissement(cantons, coso)
    equipe = econometrie.qui_est_equipe(cantons, tde, coso)
    couverture = econometrie.couverture_par_region(cantons, tde, coso)
    contre = econometrie.contrefactuel_demographique(cantons, coso)
    reconstitution = analytics.reconstitution_fri(cantons)

    ecart = perimetre.ecart_publication()

    # OBJECTIF 1 — la distance. Mesurée en mètres sur la projection UTM 31N,
    # au point représentatif de chaque canton : un centroïde tomberait hors
    # des cantons en croissant, et la distance serait fausse là où elle compte.
    rayons = accessibilite.rayons_de_marche(cantons, tde, coso)
    loin = accessibilite.deserts(cantons, tde, coso)
    voisinage = accessibilite.concentration(cantons, tde, coso)
    ensemble = voisinage[voisinage["inventaire"] == "ensemble"]

    # OBJECTIF 2 — l'état, par ses deux seuls substituts.
    maintenance = analytics.plan_de_maintenance(coso)
    service = analytics.mise_en_service(coso)

    # OBJECTIF 4 — l'indice classe-t-il autre chose que la population ?
    ampute = econometrie.fri_sans_population(cantons)

    # OBJECTIF 5 — le rattrapage, les deux programmes, et l'entretien.
    deficit = analytics.facture_rattrapage(cantons, tde, coso)
    normes = analytics.besoin_par_norme(cantons, tde, coso)
    programme = analytics.programme_ouvrages(cantons, tde, coso)
    inondations = analytics.programme_inondations(cantons, tde, coso)
    observe = analytics.taux_entretien_observe(coso)
    entretien = analytics.entretien_scenarios(
        programme, observe=observe["part_mediane"])
    budget = analytics.chaine_budgetaire(coso).set_index("etape")["montant"]
    profil = cout["profil"]["cout_par_beneficiaire"]

    geometrique = reconstitution[
        reconstitution["forme"] == "moyenne_geometrique"]

    return {
        "jeux": len(data),
        "cantons": faits["cantons"],
        "ouvrages": faits["tde_total"] + faits["coso_total"],
        "decrits": ecart["decrits"],
        "publies": ecart["communs"],
        "absents": ecart["absents"],

        "classes_cles": list(classes["classe_officielle"]),
        "classes_population": classes["population"].tolist(),
        "classes_part_equipee": matrice["part_equipee"].round(1).tolist(),
        "cantons_exposes": int(hautes["cantons"].sum()),
        "part_cantons_exposes": 100 * hautes["cantons"].sum() / faits["cantons"],
        "population_exposee": float(hautes["population"].sum()),
        "part_pop_exposee": float(hautes["part_population"].sum()),

        "fri_reproductible": int(geometrique["cantons"].iloc[0])
        if not geometrique.empty else 0,

        "regions": couverture["regions"]["region"].tolist(),
        "regions_hab_ouvrage": couverture["regions"][
            "habitants_par_ouvrage"].fillna(0).round(0).tolist(),
        "national_hab_ouvrage": couverture["national"],
        "plateaux_hab_ouvrage": float(
            couverture["regions"]["habitants_par_ouvrage"].max()),

        "r2_besoin": equipe["sans_region"]["r2"],
        "r2_region": equipe["avec_region"]["r2"],
        "prioritaires": len(prioritaires),
        "population_prioritaire": float(prioritaires["population"].sum()),

        "cout_median": cout["cout_median"],
        "cout_par_quintile": profil.round(0).tolist(),
        "ecart_cout": float(profil.iloc[0] / profil.iloc[-1]),
        "cout_beneficiaire": float(analytics.cout_par_beneficiaire(coso)
                                   ["cout_beneficiaire"].median()),

        "elasticite": float(
            elasticite["simple"]["termes"].iloc[1]["coefficient"]),
        "cantons_dotes": elasticite["cantons"],
        "gini": contre["gini"],
        "interdecile": contre["rapport_interdecile"],
        "budget_paye": float(budget["paye"]),
        "ecart_appel_offres": 100 * (
            1 - budget["contracte"] / budget["estime"]),

        "rayons_part": rayons["part"].round(1).tolist(),
        "rayons_km": rayons["rayon_km"].tolist(),
        "distance_mediane": loin["mediane_km"],
        "cantons_loin": len(loin["cantons"]),
        "seuil_loin": loin["seuil_km"],
        "part_pop_loin": loin["part_population"],
        "clark_evans": float(ensemble["R"].iloc[0]),

        "maintenance_regions": maintenance["region"].tolist(),
        "maintenance_parts": maintenance["part"].round(1).tolist(),
        "en_attente": service["en_attente"],
        "delai_remise": service["delai_median"],
        "receptionnes": service["receptionnes"],

        "rho_ampute": ampute["rho_ampute"],
        "deplacement_median": ampute["deplacement_median"],

        "manquants_mediane": deficit["ouvrages"],
        "cout_mediane": deficit["total"],
        "cout_unitaire": deficit["unitaire"],
        "observations_cout": deficit["observations"],
        "scenarios_normes": normes["scenarios"]["norme"].tolist(),
        "scenarios_manquants": normes["scenarios"]["manquants"].tolist(),
        "scenarios_cout": normes["scenarios"]["cout"].tolist(),

        "programme_annees": programme["annees"]["annee"].tolist(),
        "programme_montants": (
            programme["annees"]["montant"] / 1e9).round(2).tolist(),
        "programme_total": programme["total"],
        "programme_ouvrages": programme["ouvrages"],
        "programme_horizon": programme["horizon"],
        "programme_norme": programme["norme"],

        "inondation_total": inondations["total"],
        "inondation_cantons": inondations["cantons"],
        "inondation_horizon": inondations["horizon"],
        "inondation_unitaire": inondations["unitaire"],

        "entretien_taux_observe": 100 * observe["part_mediane"],
        "entretien_part_observee": float(
            entretien["part_investissement"].iloc[0]),
        "entretien_part_haute": float(
            entretien["part_investissement"].iloc[-1]),
    }


if __name__ == "__main__":
    demandees = sys.argv[1:] or list(LANGUES)
    partages = collecter()

    for code, chemin, pages in generer_toutes(
        PAGES, partages, langues=demandees, dossier=RACINE / "rapport"
    ):
        print(f"[{code}] {pages} pages → {chemin}")
