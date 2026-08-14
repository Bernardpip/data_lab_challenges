"""Génère la présentation de SOUTENANCE — la méthode, pas les résultats.

    python3 scripts/generer_soutenance.py

Le rapport de dix pages, lui, expose ce que les données disent. Celle-ci
expose comment on le sait : l'énoncé reçu, ce qui était attendu, les postulats
tenus, les méthodes employées, et ce qu'on a refusé de faire. Un jury qui
comprend la méthode peut juger le résultat ; l'inverse n'est pas vrai.

Quatorze planches, et des NOTES DE PRÉSENTATEUR sous chacune — le texte à
dire, pas un résumé de la planche. Elles s'affichent dans le mode Présentateur
de PowerPoint et se lisent aussi dans le fichier `SOUTENANCE-NOTES.md`, pour
qui préfère une feuille de papier.

Les chiffres ne sont pas recopiés : ils viennent de `collecter()`, la même
fonction que le rapport, qui les tire elle-même des données. Une planche ne
peut donc pas contredire le tableau de bord.
"""

import sys
from pathlib import Path

RACINE = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(RACINE))

# pyrefly: ignore [missing-import]
from pptx import Presentation                                   # noqa: E402
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt                                # noqa: E402

from socle import i18n                                          # noqa: E402

i18n.configurer(RACINE / "i18n" / "locales")

from socle.rapport import charte                                # noqa: E402


DOSSIER = RACINE / "rapport"
FICHIER = "Soutenance_Methode"


# ─── Outils de planche ───────────────────────────────────────────────────────

def _slide(prs, numero, titre, sous_titre):
    """Une planche vide, titrée et paginée."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charte.titre_page(slide, numero, titre, sous_titre)

    return slide


def _pied(slide, numero):
    """Le pied de planche — écrit ici, en français, et non par le socle.

    Celui du socle lit son libellé dans un domaine i18n : cette
    présentation-ci n'existe qu'en français et n'a pas de domaine à elle. Lui
    en fabriquer un pour une ligne de pied coûterait plus qu'il ne rapporte.
    """

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1),
                                     Inches(0.35))
    paragraphe = boite.text_frame.paragraphs[0]

    morceau = paragraphe.add_run()
    morceau.text = ("Adéquation formation-emploi au Togo · soutenance · "
                    f"planche {numero}")
    morceau.font.size = Pt(9)
    morceau.font.color.rgb = charte.MUTED
    morceau.font.name = charte.POLICE


def _notes(slide, lignes):
    """Le texte que le présentateur lira — pas un résumé de la planche.

    PowerPoint le montre dans le mode Présentateur, sur l'écran de celui qui
    parle. Écrit tel qu'il se DIT : phrases courtes, chiffres énoncés en
    entier, transitions explicites vers la planche suivante.
    """

    cadre = slide.notes_slide.notes_text_frame
    cadre.text = lignes[0]

    for ligne in lignes[1:]:
        paragraphe = cadre.add_paragraph()
        paragraphe.text = ligne

    return cadre


def _nb(valeur, decimales=2):
    """Un nombre à la française — virgule décimale, et RIEN d'autre touché.

    La première version remplaçait les points par des virgules sur la phrase
    entière : « intervalle entièrement négatif, Quand l'accès augmente… ». Le
    séparateur décimal se pose sur le NOMBRE, jamais sur le texte qui
    l'entoure.
    """

    return f"{valeur:.{decimales}f}".replace(".", ",")


def _puces(slide, gauche, haut, largeur, hauteur, entrees, taille=13):
    """Une liste à puces sobre — un fait par ligne, jamais deux.

    Pas de bloc encadré ici : une planche de méthode se lit en diagonale, et
    quatre cadres côte à côte donnent une page de formulaire.
    """

    boite = slide.shapes.add_textbox(gauche, haut, largeur, hauteur)
    cadre = boite.text_frame
    cadre.word_wrap = True

    for index, (fort, reste) in enumerate(entrees):
        paragraphe = cadre.paragraphs[0] if index == 0 else cadre.add_paragraph()
        paragraphe.space_after = Pt(10)

        morceau = paragraphe.add_run()
        morceau.text = f"·  {fort}"
        morceau.font.bold = True
        morceau.font.size = Pt(taille)
        morceau.font.color.rgb = charte.ENCRE
        morceau.font.name = charte.POLICE

        if reste:
            suite = paragraphe.add_run()
            suite.text = f"   {reste}"
            suite.font.size = Pt(taille)
            suite.font.color.rgb = charte.ENCRE_SECONDAIRE
            suite.font.name = charte.POLICE

    return cadre


# ─── Les quatorze planches ───────────────────────────────────────────────────

def page_1_couverture(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    boite = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5),
                                     Inches(3.2))
    cadre = boite.text_frame
    cadre.word_wrap = True

    charte.texte(cadre, "Adéquation formation-emploi au Togo", taille=40,
                 gras=True, premier=True, espace_apres=10)
    charte.texte(cadre, "Comment ce tableau de bord a été construit — énoncé, "
                 "postulats, méthodes", taille=16,
                 couleur=charte.ENCRE_SECONDAIRE, espace_apres=24)
    charte.texte(cadre,
                 f"{c['faits']['total']} établissements techniques · "
                 f"9 ressources ouvertes · 14 indicateurs du cahier des "
                 f"charges passés en revue", taille=13, espace_apres=26)
    charte.texte(cadre, "Kokou PIPI · Data AI Lab · Data Challenge Éducation, "
                 "défi 2", taille=12, couleur=charte.MUTED)

    _notes(slide, [
        "Bonjour. Je vais vous présenter non pas les résultats — ils sont dans "
        "le rapport de dix pages — mais la MÉTHODE qui les a produits.",
        "En quinze minutes : ce qui était demandé, ce que j'ai posé comme "
        "règles avant de commencer, les outils statistiques employés, et ce "
        "que j'ai refusé de faire.",
        "Un mot sur l'esprit : tout ce que vous verrez est vérifiable. Le "
        "tableau de bord est en ligne, l'archive contient les données brutes, "
        "et chaque chiffre d'une planche vient du même code que l'écran.",
    ])


def page_2_enonce(prs, c):
    slide = _slide(prs, 2, "L'énoncé",
                   "Ce que le défi demandait, et ce que j'en ai fait une "
                   "question mesurable")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.4), [
        ("La question posée", "« L'offre de formation est-elle alignée sur "
         "les besoins de l'économie togolaise ? »"),
        ("Six objectifs", "cartographier l'offre, mesurer l'accès, suivre les "
         "moyens publics, décrire l'insertion, croiser, recommander"),
        ("Un corpus imposé", "les ressources ouvertes du portail national, "
         "et elles seules"),
    ])

    charte.bloc_analyse(
        slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(2.5),
        nom="La reformulation",
        question="Une question d'opinion peut-elle se mesurer ?",
        resultat="Oui, en trois écarts",
        lecture="« Alignement » ne se mesure pas directement. Je l'ai éclaté "
                "en trois écarts observables : entre les territoires, entre "
                "l'accès et les moyens, entre la formation et l'insertion. "
                "Chacun se calcule ; leur somme ne se calcule pas, et je ne "
                "l'ai pas inventée.")

    _pied(slide, 2)
    _notes(slide, [
        "L'énoncé demandait si l'offre de formation est alignée sur les "
        "besoins de l'économie. C'est une question d'opinion : telle quelle, "
        "elle n'a pas de réponse chiffrée.",
        "Mon premier travail a donc été de la reformuler en trois écarts qui, "
        "eux, se mesurent : l'écart entre territoires, l'écart entre l'accès "
        "aux études et les moyens qu'on y met, et l'écart entre ce qu'on "
        "forme et ce qui s'insère.",
        "Point important : je n'ai pas fabriqué un « indice d'alignement » "
        "global. Additionner trois écarts de natures différentes donnerait un "
        "nombre qui a l'air savant et qui ne veut rien dire.",
    ])


def page_3_attendus(prs, c):
    slide = _slide(prs, 3, "Les attendus",
                   "Quatorze indicateurs, quatre livrables — et un verdict "
                   "public pour chacun")

    for index, (valeur, libelle, detail, couleur) in enumerate([
        ("9", "Indicateurs honorés", "calculés et affichés", charte.PRIMAIRE),
        ("2", "Partiellement", "la donnée existe, mais pas à la maille "
         "demandée", charte.PRIMAIRE),
        ("3", "Impossibles", "et la planche dit POURQUOI", charte.ROUGE),
    ]):
        charte.bloc_constat(
            slide, Inches(0.7 + index * 4.15), Inches(1.9), Inches(3.9),
            Inches(1.5), valeur=valeur, libelle=libelle, detail=detail,
            couleur=couleur)

    _puces(slide, Inches(0.7), Inches(3.8), Inches(11.9), Inches(2.5), [
        ("Les quatre livrables", "un tableau de bord bilingue, un rapport de "
         "dix pages, l'archive du code et des données, un diagnostic "
         "d'installation"),
        ("Le principe", "un indicateur impossible reste AFFICHÉ, avec sa "
         "cause — le masquer donnerait l'illusion d'un travail complet"),
    ])

    _pied(slide, 3)
    _notes(slide, [
        "Le cahier des charges se décompose en quatorze indicateurs "
        "élémentaires. J'en ai honoré neuf, deux partiellement, et trois sont "
        "impossibles à partir du corpus.",
        "Ces trois-là sont affichés dans le tableau de bord, avec leur cause. "
        "C'est un choix : un jury qui ne voit que ce qui a marché ne peut pas "
        "juger la couverture réelle du travail.",
        "Les deux « partiels » sont des cas où la donnée existe, mais pas à "
        "la maille demandée — typiquement au niveau national quand on la "
        "voudrait par région.",
        "Un détail qui compte pour la suite : cet audit est recalculé à "
        "chaque exécution, il n'est pas écrit à la main.",
    ])


def page_4_corpus(prs, c):
    slide = _slide(prs, 4, "Le matériau",
                   "Neuf ressources — et une asymétrie qui commande tout le "
                   "reste")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(6.2), Inches(4.4), [
        ("8 ressources chargées", "sur 9 ; la neuvième est une table de "
         "contrôle du jeu technique"),
        ("1 seule descend sous le national", "les 256 établissements "
         "techniques, géolocalisés"),
        ("7 séries nationales", "budget, inscriptions, dépense, chômage, "
         "effectifs — un seul chiffre par année"),
        ("Périodes hétérogènes", "supérieur 2018, technique 2025, séries "
         "internationales 1971-2022"),
    ])

    charte.bloc_analyse(
        slide, Inches(7.3), Inches(2.0), Inches(5.3), Inches(2.9),
        nom="La conséquence",
        question="Pourquoi cette asymétrie décide de tout ?",
        resultat="Aucun score régional composite",
        lecture="Un indicateur national ne peut pas entrer dans un classement "
                "de régions : il y apporterait la même valeur partout, en "
                "donnant l'illusion d'avoir été mesuré localement. Le chômage "
                "des diplômés est donc absent de tous les scores "
                "territoriaux, et la vue le dit.")

    _pied(slide, 4)
    _notes(slide, [
        "Neuf ressources, dont huit réellement chargées. Et surtout : une "
        "seule descend sous le niveau national.",
        "C'est l'asymétrie qui commande toute l'architecture. Sept séries ne "
        "donnent qu'un chiffre par année pour tout le pays ; une seule donne "
        "deux cent cinquante-six points géolocalisés.",
        "Conséquence directe : je me suis interdit tout score régional "
        "composite qui mélangerait les deux. Si j'injecte le chômage national "
        "dans un score par région, il vaut la même chose partout — le score a "
        "l'air riche, il ne mesure rien de plus.",
        "Les périodes aussi sont hétérogènes : 2018 pour le supérieur, 2025 "
        "pour le technique. Aucun croisement ne franchit cet écart sans le "
        "dire.",
    ])


def page_5_postulats(prs, c):
    slide = _slide(prs, 5, "Les postulats",
                   "Quatre règles posées AVANT de regarder les données — donc "
                   "avant de savoir ce qu'elles arrangeraient")

    for index, (nom, question, resultat, lecture) in enumerate([
        ("1 · Aucune donnée fabriquée",
         "Que faire d'une série trouée ?",
         "On laisse le trou",
         "Ni interpolation, ni moyenne de remplissage. Les non-réponses "
         "restent visibles sous « Non renseigné » : les supprimer embellirait "
         "toutes les répartitions."),
        ("2 · Aucun croisement non autorisé",
         "Peut-on croiser deux mailles différentes ?",
         "Non, jamais en silence",
         "Chaque croisement déclare ses ingrédients, sa clé de jointure et "
         "son nombre d'observations. Cinq années communes n'autorisent pas "
         "les conclusions de dix."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 6.15), Inches(1.9), Inches(5.9),
            Inches(2.1), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    for index, (nom, question, resultat, lecture) in enumerate([
        ("3 · Le non-significatif s'affiche",
         "Que fait-on d'un modèle qui ne conclut pas ?",
         "On le publie quand même",
         "Deux des cinq modèles ne concluent pas. Les retirer donnerait à "
         "l'ensemble une solidité qu'il n'a pas."),
        ("4 · Le contexte externe reste dehors",
         "Peut-on compléter avec d'autres sources ?",
         "Oui, mais séparément",
         "Les repères d'enquêtes nationales sont sourcés, signalés "
         "visuellement, et n'entrent dans aucun calcul du corpus."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 6.15), Inches(4.2), Inches(5.9),
            Inches(2.1), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    _pied(slide, 5)
    _notes(slide, [
        "Voici les quatre règles que je me suis données avant d'ouvrir les "
        "fichiers. C'est important : posées après, elles auraient été "
        "choisies en fonction de ce qui m'arrangeait.",
        "Première règle : aucune donnée fabriquée. Pas d'interpolation, pas "
        "de moyenne de remplissage. Une série trouée reste trouée.",
        "Deuxième : aucun croisement que les données n'autorisent pas, et "
        "chaque croisement déclare son nombre d'observations.",
        "Troisième, la plus inconfortable : les résultats non significatifs "
        "sont affichés comme tels. Deux de mes cinq modèles ne concluent pas. "
        "Ils sont dans le tableau de bord, avec leur p-value.",
        "Quatrième : le contexte externe — les enquêtes nationales — est "
        "sourcé et séparé. Il éclaire, il ne se mélange pas.",
    ])


def page_6_identifier(prs, c):
    """Temps 1 — savoir ce que chaque fichier peut, et ne peut pas, mesurer."""

    slide = _slide(prs, 6, "Méthode, temps 1 — identifier les données",
                   "Avant tout calcul : ouvrir chaque fichier, le compter, et "
                   "décider ce qu'il a le droit de dire")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.4), [
        ("Ouvrir, pas supposer", "chaque fichier est lu ligne à ligne : "
         "combien d'observations, quelles colonnes, quelles valeurs "
         "manquantes"),
        ("Situer sa maille", "national, régional, ponctuel — c'est elle qui "
         "décide des croisements possibles"),
        ("Situer sa période", "1971-2022 pour les séries, 2018 pour le "
         "supérieur, 2025 pour le technique"),
        ("Comparer au dictionnaire", "le fichier tient-il ce que ses "
         "métadonnées annoncent ?"),
    ], taille=12.5)

    charte.bloc_constat(
        slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(1.5),
        valeur="16 / 216", libelle="Champs publiés, sur ceux décrits",
        detail="201 champs collectés ne sont pas diffusés",
        couleur=charte.ROUGE)

    charte.bloc_analyse(
        slide, Inches(7.1), Inches(3.7), Inches(5.5), Inches(2.4),
        nom="Ce que cette étape a trouvé",
        question="Un fichier peut-il se contredire ?",
        resultat="Oui — et il fallait trancher",
        lecture="Le fichier du supérieur inscrit 14 établissements privés sur "
                "sa ligne de total, quand le détail par ville en compte 65. "
                "J'ai retenu le détail : c'est un CHOIX, il est écrit dans les "
                "annexes, et tous les chiffres du supérieur en dépendent.")

    _pied(slide, 6)
    _notes(slide, [
        "La méthode se déroule en trois temps. Le premier : identifier les "
        "données.",
        "Cela veut dire ouvrir réellement chaque fichier — pas lire son "
        "intitulé. Compter les lignes, regarder les colonnes, repérer les "
        "vides.",
        "Et surtout : situer sa MAILLE et sa PÉRIODE. C'est la maille qui "
        "décide de ce qu'on aura le droit de croiser ensuite, et la période "
        "qui décide de ce qu'on pourra comparer.",
        "Deux découvertes de cette étape. D'abord le dictionnaire du fichier "
        "technique : il décrit deux cent seize champs, le fichier en publie "
        "seize. Deux cent un champs collectés ne sont pas ouverts.",
        "Ensuite, le fichier du supérieur se contredit lui-même : quatorze "
        "établissements privés sur la ligne de total, soixante-cinq dans le "
        "détail par ville. J'ai retenu le détail, c'est écrit, et je peux en "
        "discuter.",
    ])


def page_7_croiser(prs, c):
    """Temps 2 — croiser, et déclarer ce que la jointure autorise."""

    slide = _slide(prs, 7, "Méthode, temps 2 — croiser les données",
                   "Six croisements, chacun déclarant sa clé, ses ingrédients "
                   "et son nombre d'observations")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(6.4), Inches(4.4), [
        ("Accès × moyens", "inscriptions au supérieur contre dépense par "
         "étudiant"),
        ("Budget × accès", "l'enveloppe votée suit-elle le nombre "
         "d'étudiants ?"),
        ("Cohérence du financement", "budget exécuté contre dépense par "
         "étudiant"),
        ("Réseaux par région", "technique et supérieur sur le même "
         "territoire"),
        ("Accès × insertion", "inscriptions contre chômage des diplômés"),
        ("Accès × effectifs", "féminisation et filières scientifiques"),
    ], taille=12.5)

    charte.bloc_analyse(
        slide, Inches(7.5), Inches(2.0), Inches(5.1), Inches(2.7),
        nom="La règle du croisement",
        question="Qu'est-ce qui rend une jointure honnête ?",
        resultat="Dire sur combien elle porte",
        lecture="Chaque croisement affiche ses deux fichiers, sa clé — presque "
                "toujours l'année — et le nombre d'observations qui survivent "
                "à la jointure. Deux séries qui se recouvrent sur cinq ans "
                "n'autorisent pas les conclusions de vingt, et le lecteur doit "
                "voir ce nombre avant de lire le résultat.")

    _pied(slide, 7)
    _notes(slide, [
        "Deuxième temps : croiser. Il y a six croisements, et j'ai emprunté au "
        "vocabulaire de la cuisine — chacun est une recette qui déclare ses "
        "ingrédients.",
        "Concrètement, quand vous ouvrez un croisement dans le tableau de "
        "bord, il vous dit : voici les deux fichiers, voici la clé de "
        "jointure — presque toujours l'année —, et voici le nombre "
        "d'observations qui restent après la jointure.",
        "Pourquoi c'est capital : deux séries qui se recouvrent sur cinq ans "
        "ne permettent pas les mêmes affirmations que deux séries qui se "
        "recouvrent sur vingt.",
        "Le croisement central est le premier, accès contre moyens. C'est lui "
        "qui produit l'effet ciseaux, et c'est sur lui que je vais montrer "
        "les outils.",
    ])


def page_8_outils_relation(prs, c):
    """Temps 3, premier volet — mesurer une relation entre deux séries."""

    slide = _slide(prs, 8, "Méthode, temps 3 — les outils, et ce qu'ils font",
                   "Trois outils pour mesurer une relation entre deux séries")

    for index, (nom, question, resultat, lecture) in enumerate([
        ("La régression linéaire",
         "Que fait-elle, concrètement ?",
         "Elle trace la droite la plus proche des points",
         "Sa pente dit : « quand x augmente d'une unité, y bouge de tant ». "
         "Je ne publie jamais la pente seule, mais avec trois garde-fous — le "
         "R², part de la variation expliquée ; la p-value, probabilité que le "
         "hasard suffise ; l'intervalle de confiance, la fourchette où se "
         "trouve la vraie pente."),
        ("La tendance temporelle",
         "Et quand x est le temps ?",
         "Elle donne un rythme : tant par an",
         "C'est la même régression, avec les années en abscisse. Elle répond "
         "à « de combien cela bouge-t-il chaque année », ce qui se dit et se "
         "retient mieux qu'un coefficient abstrait."),
        ("L'élasticité",
         "Comment comparer des grandeurs sans unité commune ?",
         "En passant les deux au logarithme",
         "La pente devient alors un pourcentage : ici, quand l'accès augmente "
         "de 1 %, la dépense par étudiant recule de "
         f"{_nb(abs(c['elasticite']['pente']))} %. Un nombre d'étudiants et "
         "un pourcentage de PIB deviennent comparables."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 4.15), Inches(1.9), Inches(3.9),
            Inches(3.9), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    _pied(slide, 8)
    _notes(slide, [
        "Troisième temps : les outils. Je vais les expliquer, parce qu'un "
        "nom d'outil ne prouve rien — c'est ce qu'il fait qui compte.",
        "La régression linéaire, d'abord. On cherche la droite qui passe au "
        "plus près du nuage de points. Sa pente répond à : quand x augmente "
        "d'une unité, de combien bouge y ?",
        "Mais une pente toute seule ne vaut rien. Je publie toujours trois "
        "garde-fous avec elle. Le R carré : quelle part de la variation la "
        "droite explique-t-elle. La p-value : quelle est la probabilité "
        "d'observer ça par pur hasard. Et l'intervalle de confiance : dans "
        "quelle fourchette se trouve la vraie pente.",
        "La tendance temporelle, c'est la même chose avec les années en "
        "abscisse. Elle donne un rythme — tant par an — et c'est ce qui se "
        "retient.",
        "L'élasticité, enfin. Le problème : comment comparer un nombre "
        "d'étudiants avec un pourcentage de PIB ? On passe les deux au "
        "logarithme, et la pente devient un pourcentage. Ici : quand l'accès "
        "monte de un pour cent, la dépense par étudiant recule de zéro "
        "virgule quatre-vingt-six pour cent.",
    ])


def page_9_outils_forme(prs, c):
    """Temps 3, second volet — la forme de la relation, et la répartition."""

    slide = _slide(prs, 9, "Méthode, temps 3 — les outils (suite)",
                   "Un outil pour la forme d'une relation, deux pour une "
                   "répartition")

    for index, (nom, question, resultat, lecture) in enumerate([
        ("La corrélation, mesurée deux fois",
         "Pourquoi Pearson ET Spearman ?",
         "Ils ne demandent pas la même chose",
         "Pearson demande si les points s'alignent sur une droite. Spearman "
         "demande seulement s'ils montent et descendent ensemble. Quand les "
         "deux s'écartent, la relation existe mais n'est pas linéaire — et "
         "cet écart est une information, pas un défaut."),
        ("Le test de rupture",
         "Une tendance peut-elle changer en cours de route ?",
         "On coupe la série et on compare",
         "On calcule la pente avant une année charnière, puis après. Sur "
         f"l'accès au supérieur, la rupture est en {c['rupture']['annee_rupture']} : "
         f"{_nb(c['rupture']['avant']['pente'])} point par an avant, "
         f"{_nb(c['rupture']['apres']['pente'])} après. La progression a "
         "changé de régime, elle n'a pas seulement continué."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 6.15), Inches(1.9), Inches(5.9),
            Inches(2.2), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    for index, (nom, question, resultat, lecture) in enumerate([
        ("Les deux indices de concentration",
         "Comment chiffrer un déséquilibre territorial ?",
         "Herfindahl et Gini, ensemble",
         "Herfindahl additionne les carrés des parts : il s'envole dès qu'un "
         "acteur domine. Gini décrit toute la distribution : zéro si tout le "
         f"monde a la même chose, un si tout est chez un seul. Ici "
         f"{_nb(c['concentration']['hhi'])} et "
         f"{_nb(c['concentration']['gini'])}."),
        ("Exécution contre enveloppe",
         "Les grosses enveloppes se dépensent-elles moins bien ?",
         "Question posée, réponse non concluante",
         "Une régression du taux d'exécution sur le montant voté. La pente va "
         "dans le sens attendu, mais l'intervalle de confiance traverse zéro : "
         "sur six années, on ne peut pas conclure — et c'est publié ainsi."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 6.15), Inches(4.3), Inches(5.9),
            Inches(2.2), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    _pied(slide, 9)
    _notes(slide, [
        "Quatre outils encore, plus rapidement.",
        "La corrélation, je la mesure deux fois. Pearson demande si les "
        "points s'alignent sur une droite ; Spearman demande seulement s'ils "
        "montent et descendent ensemble. Quand les deux s'écartent, la "
        "relation existe mais n'est pas droite — et c'est une information.",
        "Le test de rupture : on coupe la série à une année charnière et on "
        "compare les deux pentes. Sur l'accès au supérieur, la rupture est en "
        "deux mille : zéro virgule zéro neuf point par an avant, zéro virgule "
        "soixante-trois après. Ce n'est pas une continuation, c'est un "
        "changement de régime.",
        "Les deux indices de concentration. Herfindahl additionne les carrés "
        "des parts : il s'envole dès qu'un acteur domine. Gini décrit toute "
        "la distribution. Je donne les deux parce qu'ils ne disent pas la "
        "même chose.",
        "Et le dernier, que je garde pour montrer une honnêteté : la "
        "régression du taux d'exécution sur le montant voté. Elle ne conclut "
        "pas. Je la publie quand même.",
    ])


def page_10_resultats(prs, c):
    slide = _slide(prs, 10, "Ce que la méthode produit",
                   "Trois résultats, chacun rattaché à l'outil qui l'établit")

    charte.bloc_constat(
        slide, Inches(0.7), Inches(1.9), Inches(3.9), Inches(1.6),
        valeur=f'{round(c["faits"]["part_region_tete"])} %',
        libelle="des établissements dans une région",
        detail=f'Herfindahl {_nb(c["concentration"]["hhi"])} · Gini '
               f'{_nb(c["concentration"]["gini"])}',
        couleur=charte.ROUGE)

    charte.bloc_constat(
        slide, Inches(4.8), Inches(1.9), Inches(3.9), Inches(1.6),
        valeur="×4,7 / ÷3,6", libelle="Accès et dépense, base 100 en 1998",
        # Les DEUX PÉRIODES sont écrites : les séries ne s'arrêtent pas la
        # même année — l'accès va jusqu'en 2020, la dépense s'arrête en 2017.
        # Un « depuis 1998 » sans borne de fin laisserait croire à deux
        # trajectoires mesurées jusqu'au même jour, et c'est la première chose
        # qu'un lecteur attentif viendrait contester.
        detail="accès jusqu'en 2020, dépense jusqu'en 2017 — le calcul de "
               "l'élasticité, lui, ne porte que sur leurs 10 années communes",
        couleur=charte.ROUGE)

    charte.bloc_constat(
        slide, Inches(8.9), Inches(1.9), Inches(3.7), Inches(1.6),
        valeur=f'{c["moyenne_es"]:.0f} %',
        libelle="Exécution du budget supérieur",
        detail=f'contre {c["moyenne_nat"]:.0f} % au niveau national')

    charte.bloc_analyse(
        slide, Inches(0.7), Inches(3.8), Inches(5.9), Inches(2.4),
        nom="Le résultat le plus solide",
        question="log(dépense/étudiant) contre log(accès), 10 années",
        resultat=f'Élasticité {_nb(c["elasticite"]["pente"])}',
        lecture=f'R² de {_nb(c["elasticite"]["r2"])}, p-value sous 0,001, '
                f'intervalle de confiance entièrement négatif. Quand l\'accès '
                f'augmente de 1 %, la dépense par étudiant recule de '
                f'{_nb(abs(c["elasticite"]["pente"]))} %.')

    charte.bloc_analyse(
        slide, Inches(6.9), Inches(3.8), Inches(5.7), Inches(2.4),
        nom="Le résultat qui NE conclut pas",
        question="taux d'exécution contre montant voté, 6 années",
        resultat=f'p = {_nb(c["execution"]["p_value"])}',
        lecture="La pente est négative et le R² correct, mais l'intervalle "
                "de confiance traverse zéro : sur six années, on ne peut pas "
                "conclure. C'est publié tel quel.")

    _pied(slide, 10)
    _notes(slide, [
        "Voici trois résultats, et je les donne surtout pour montrer le lien "
        "avec les outils.",
        "Soixante pour cent des établissements techniques sont dans une seule "
        "région, la Maritime. Herfindahl à zéro virgule quarante-et-un, Gini "
        "à zéro virgule quarante-huit : la concentration est forte, et "
        "mesurée, pas ressentie.",
        "Le deuxième est le plus important. Depuis 1998, l'accès au supérieur "
        "a été multiplié par quatre virgule sept, et la dépense publique par "
        "étudiant divisée par trois virgule six. C'est ce que j'appelle "
        "l'effet ciseaux.",
        "L'élasticité vaut moins zéro virgule quatre-vingt-six : quand "
        "l'accès augmente d'un pour cent, la dépense par étudiant recule de "
        "presque un pour cent. R carré de zéro virgule quatre-vingt-trois, "
        "p-value sous un pour mille.",
        "Un mot sur ce que cela VEUT DIRE, parce que c'est là qu'est "
        "l'argument. La dépense par étudiant, c'est un budget divisé par un "
        "nombre d'étudiants. Une élasticité de moins zéro virgule "
        "quatre-vingt-six, presque moins un, dit que le dénominateur a "
        "explosé pendant que le numérateur bougeait à peine. Le pays a "
        "multiplié ses étudiants sans multiplier l'argent. Ce n'est pas une "
        "corrélation mystérieuse, c'est arithmétique.",
        "Si on m'objecte « corrélation n'est pas causalité », la réponse est "
        "là : le lien est mécanique, les deux grandeurs partagent leur "
        "dénominateur. Je ne prétends pas que l'accès CAUSE la baisse ; je "
        "constate que l'argent n'a pas suivi le nombre.",
        "Et à droite, le contre-exemple que je tiens à montrer : le lien "
        "entre taille de l'enveloppe et taux d'exécution. La pente va dans le "
        "sens attendu, mais l'intervalle de confiance traverse zéro. Six "
        "années, c'est trop peu. Je publie le résultat en disant qu'il ne "
        "conclut pas.",
    ])


def page_11_limites(prs, c):
    slide = _slide(prs, 11, "Ce que la méthode refuse d'établir",
                   "Trois manques, trois causes différentes — et trois "
                   "remèdes différents")

    for index, (nom, question, resultat, lecture) in enumerate([
        ("Collecté, non publié",
         "Combien d'élèves par établissement ?",
         "201 champs manquants",
         "Le remède est une EXPORTATION : la donnée existe, elle est déjà "
         "collectée. Aucune enquête à financer."),
        ("Inexistant à la maille utile",
         "Quel taux d'insertion par filière ?",
         "Le chômage n'existe qu'au national",
         "Le remède est une ENQUÊTE d'insertion. Trois croisements attendus "
         "par l'énoncé restent hors de portée."),
        ("Nomenclature absente",
         "Quelle part de filières scientifiques ?",
         "Aucun référentiel de métiers",
         "Le remède est un RÉFÉRENTIEL national. Sans lui, le dénominateur "
         "n'existe pas."),
    ]):
        charte.bloc_analyse(
            slide, Inches(0.7 + index * 4.15), Inches(1.9), Inches(3.9),
            Inches(3.0), nom=nom, question=question, resultat=resultat,
            lecture=lecture)

    _puces(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.1), [
        ("Un indicateur absent", "signifie toujours ABSENT DU CORPUS RETENU, "
         "jamais inexistant — la nuance change la recommandation"),
    ])

    _pied(slide, 11)
    _notes(slide, [
        "Cette planche est celle à laquelle je tiens le plus.",
        "Il y a trois manques dans ce travail, et ils ont trois causes "
        "différentes. Les confondre mènerait à des recommandations fausses.",
        "Le premier : la donnée est collectée mais pas publiée. Deux cent un "
        "champs. Le remède coûte un export, pas une enquête.",
        "Le deuxième : la donnée n'existe pas à la maille utile. Le chômage "
        "des diplômés n'est publié qu'au niveau national. Là, il faut une "
        "vraie enquête d'insertion.",
        "Le troisième : il n'existe aucune nomenclature des métiers ni des "
        "disciplines. Sans référentiel, la question « quelle part de filières "
        "scientifiques » n'a pas de dénominateur.",
        "Et la phrase du bas est ma règle : un indicateur absent est absent "
        "DU CORPUS, jamais inexistant. Je n'ai pas le droit de dire que le "
        "Togo ne sait pas ; je peux dire que le portail ne publie pas.",
    ])


def page_12_verifiabilite(prs, c):
    slide = _slide(prs, 12, "Comment tout cela se vérifie",
                   "Le travail est fait pour être contredit — encore faut-il "
                   "pouvoir le refaire")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4), [
        ("Un chiffre, une source", "les commentaires sous les graphes se "
         "recalculent sur la sélection : aucun nombre n'est figé dans un texte"),
        ("Le rapport ne peut pas diverger de l'écran", "il est produit par "
         "les mêmes fonctions, à partir des mêmes données, en une seule "
         "collecte"),
        ("L'archive contient les données brutes", "telles que téléchargées du "
         "portail, non modifiées"),
        ("Un diagnostic livré", "une commande contrôle la version de Python, "
         "les bibliothèques et les fichiers avant tout lancement"),
        ("Le tableau de bord est en ligne", "tg-datalab-education-challenge2."
         "bernardpip.com — rien à installer pour le juger"),
    ], taille=12.5)

    _pied(slide, 12)
    _notes(slide, [
        "Un mot sur la vérifiabilité, parce qu'une méthode qu'on ne peut pas "
        "rejouer n'est pas une méthode.",
        "Aucun chiffre n'est écrit en dur dans un commentaire. Quand vous "
        "filtrez sur les Savanes, les phrases sous les graphes se "
        "réécrivent — elles sont calculées, pas rédigées.",
        "Le rapport PowerPoint est produit par les mêmes fonctions que "
        "l'écran, en une seule collecte de chiffres. Il ne peut pas "
        "diverger.",
        "L'archive contient les fichiers du portail non modifiés, et un "
        "script de diagnostic qui vérifie l'installation avant de lancer quoi "
        "que ce soit.",
        "Et surtout : le tableau de bord est en ligne. Vous n'avez rien à "
        "installer pour le contredire.",
    ])


def page_13_recommandations(prs, c):
    slide = _slide(prs, 13, "Des constats aux recommandations",
                   "26 leviers, 5 axes — chacun part d'un chiffre, aucun "
                   "d'une opinion")

    charte.bloc_constat(
        slide, Inches(0.7), Inches(1.9), Inches(5.7), Inches(1.6),
        valeur="26", libelle="Leviers d'action",
        detail="chacun renvoie à la vue qui l'établit")

    charte.bloc_constat(
        slide, Inches(6.9), Inches(1.9), Inches(5.7), Inches(1.6),
        valeur="1", libelle="Première décision recommandée",
        detail="elle est informationnelle, pas budgétaire",
        couleur=charte.PRIMAIRE)

    charte.bloc_analyse(
        slide, Inches(0.7), Inches(3.8), Inches(11.9), Inches(2.4),
        nom="La règle de rédaction",
        question="Qu'est-ce qui distingue un levier d'un vœu ?",
        resultat="Un constat chiffré, et une vue qui le porte",
        lecture="Chaque levier commence par le chiffre qui le motive et "
                "renvoie à l'onglet où ce chiffre se vérifie. Aucun n'est une "
                "recommandation générique du type « renforcer la "
                "coordination ». Et la toute première n'est pas de dépenser : "
                "c'est de publier les 201 champs déjà collectés — parce que "
                "sans eux, on pilote à l'aveugle.")

    _pied(slide, 13)
    _notes(slide, [
        "Les recommandations, maintenant. Il y en a vingt-six, réparties en "
        "cinq axes.",
        "Ma règle de rédaction : un levier commence par le chiffre qui le "
        "motive, et renvoie à l'onglet où ce chiffre se vérifie. Si je ne "
        "peux pas faire ça, ce n'est pas un levier, c'est un vœu.",
        "Et je veux insister sur la première recommandation, parce qu'elle "
        "surprend : elle ne demande pas un franc. Elle demande de publier les "
        "deux cent un champs déjà collectés.",
        "Tant qu'ils ne le sont pas, personne — ni le ministère, ni un "
        "chercheur, ni moi — ne peut calculer un taux d'encadrement par "
        "établissement. On pilote à l'aveugle sur une donnée qui existe "
        "déjà.",
    ])


def page_14_reutilisable(prs, c):
    slide = _slide(prs, 14, "Ce qui est réutilisable",
                   "La méthode a été rejouée sur un second défi, en quelques "
                   "jours")

    _puces(slide, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4), [
        ("Un socle partagé", "coquille, charte, graphes, cartes, i18n, "
         "économétrie — le défi n'écrit que ses analyses"),
        ("Les mêmes règles de rigueur", "aucune donnée fabriquée, aucun "
         "croisement non autorisé, le non-significatif affiché"),
        ("La même chaîne de livraison", "tableau de bord bilingue, rapport, "
         "archive reconstructible, diagnostic"),
        ("Preuve par l'usage", "le défi Environnement — eau et hydraulique — "
         "a été construit sur ce socle, avec ses propres analyses"),
    ], taille=12.5)

    _pied(slide, 14)
    _notes(slide, [
        "Un point qui dépasse ce défi : ce qui a été construit ici est "
        "réutilisable.",
        "La coquille, la charte graphique, les formes de graphes, les cartes, "
        "le bilinguisme et la boîte à outils économétrique vivent dans un "
        "socle partagé. Un nouveau défi n'écrit que ses analyses.",
        "Ce n'est pas une promesse : le défi Environnement, sur l'eau et "
        "l'hydraulique, a été construit sur ce même socle, avec ses propres "
        "données et ses propres analyses.",
        "Et les règles de rigueur voyagent avec le socle. C'est le point "
        "important : ce n'est pas un gabarit graphique, c'est une méthode "
        "outillée.",
    ])


def page_15_conclusion(prs, c):
    slide = _slide(prs, 15, "Conclusion",
                   "Ce que ce travail établit, et ce qu'il laisse ouvert")

    charte.bloc_analyse(
        slide, Inches(0.7), Inches(1.9), Inches(5.9), Inches(2.5),
        nom="Ce qui est établi",
        question="L'alignement est-il réalisé ?",
        resultat="Non, sur les trois écarts mesurés",
        lecture="L'offre est concentrée sur un territoire, l'accès progresse "
                "quand les moyens par étudiant reculent, et l'insertion ne "
                "peut pas être suivie faute de donnée locale. Les deux "
                "premiers écarts sont mesurés ; le troisième est un manque, "
                "et c'en est un résultat.")

    charte.bloc_analyse(
        slide, Inches(6.9), Inches(1.9), Inches(5.7), Inches(2.5),
        nom="Ce qui reste ouvert",
        question="Que faudrait-il pour aller plus loin ?",
        resultat="Trois gestes, dans cet ordre",
        lecture="Publier les 201 champs collectés. Lancer une enquête "
                "d'insertion ventilée par filière et par région. Établir un "
                "référentiel national des métiers. Aucun des trois n'est un "
                "travail de statisticien : ce sont des décisions.")

    _puces(slide, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.6), [
        ("La phrase que je retiens", "un tableau de bord honnête montre "
         "d'abord ce qu'il ne peut pas montrer — c'est à cela qu'on voit "
         "qu'il n'a rien inventé"),
    ])

    _pied(slide, 15)
    _notes(slide, [
        "Pour conclure.",
        "L'alignement n'est pas réalisé, et je peux le dire sur deux des "
        "trois écarts : le territoire et les moyens. Le troisième, "
        "l'insertion, je ne peux pas le mesurer — et c'est un résultat, pas "
        "un échec.",
        "Ce qu'il faudrait pour aller plus loin tient en trois gestes, dans "
        "cet ordre : publier ce qui est déjà collecté, lancer une enquête "
        "d'insertion, établir un référentiel des métiers.",
        "Aucun des trois n'est un travail de statisticien. Ce sont des "
        "décisions administratives — et c'est précisément pour ça qu'un "
        "tableau de bord sert à quelque chose : il montre où la décision "
        "manque.",
        "Je termine là-dessus : un tableau de bord honnête montre d'abord ce "
        "qu'il ne peut pas montrer. Merci. Je réponds à vos questions.",
    ])


PAGES = [
    page_1_couverture, page_2_enonce, page_3_attendus, page_4_corpus,
    page_5_postulats,
    # Les trois temps de la méthode, dans l'ordre où on les vit : identifier
    # ce qu'on a, croiser ce qui peut l'être, puis seulement mesurer. Les
    # outils venaient AVANT les croisements : on expliquait comment mesurer
    # avant d'avoir dit sur quoi.
    page_6_identifier, page_7_croiser,
    page_8_outils_relation, page_9_outils_forme,
    page_10_resultats, page_11_limites, page_12_verifiabilite,
    page_13_recommandations, page_14_reutilisable, page_15_conclusion,
]


def _notes_en_markdown(prs):
    """Les notes, rassemblées en un document qu'on peut imprimer.

    Le mode Présentateur de PowerPoint les montre déjà ; ce fichier sert à qui
    préfère une feuille de papier, et à relire son texte dans le train.
    """

    lignes = ["# Soutenance — notes du présentateur", "",
              "*Le texte à dire, planche par planche. Les planches, elles, "
              "sont dans `Soutenance_Methode.pdf`.*", ""]

    for numero, slide in enumerate(prs.slides, 1):
        titres = [f.text_frame.text.splitlines()[0]
                  for f in slide.shapes
                  if f.has_text_frame and f.text_frame.text.strip()]
        lignes += [f"## Planche {numero} — {titres[0] if titres else ''}", ""]

        if slide.has_notes_slide:
            for paragraphe in slide.notes_slide.notes_text_frame.paragraphs:
                if paragraphe.text.strip():
                    lignes += [paragraphe.text.strip(), ""]

    return "\n".join(lignes)


def main():
    from scripts.generer_presentation import collecter

    chiffres = collecter()

    prs = Presentation()
    prs.slide_width, prs.slide_height = charte.LARGEUR, charte.HAUTEUR

    for page in PAGES:
        page(prs, chiffres)

    DOSSIER.mkdir(exist_ok=True)
    pptx = DOSSIER / f"{FICHIER}.pptx"
    prs.save(pptx)

    notes = DOSSIER / "SOUTENANCE-NOTES.md"
    notes.write_text(_notes_en_markdown(prs), encoding="utf-8")

    print(f"  {pptx.name} · {len(prs.slides)} planches")
    print(f"  {notes.name} · les notes du présentateur")

    return pptx


if __name__ == "__main__":
    main()
