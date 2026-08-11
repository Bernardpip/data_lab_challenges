"""Génère le rapport PowerPoint (10 pages maximum) attendu comme livrable.

Lancement :
    python scripts/generer_presentation.py            # les deux langues
    python scripts/generer_presentation.py fr         # une seule

Choix techniques :
  · les graphiques sont NATIFS PowerPoint (python-pptx `add_chart`), donc
    éditables et redimensionnables par le lecteur — pas des images figées ;
  · tous les chiffres sont recalculés depuis les jeux nettoyés au moment de la
    génération : la présentation ne peut pas diverger du tableau de bord ;
  · la charte reprend les tokens du dashboard (mêmes couleurs, même hiérarchie
    typographique), pour que les deux livrables se lisent comme un seul travail.

Bilingue. Aucun texte n'est écrit ici : tout vient de `i18n/locales/
presentation.json`, comme pour les vues. Deux différences avec le tableau de
bord, imposées par le fait qu'un fichier généré n'a pas de session :

  · la langue est un ARGUMENT, jamais l'état de session. `t()` lit
    `st.session_state`, ce qui n'existe pas quand le script tourne en ligne de
    commande, et surtout ne dirait rien de la langue DEMANDÉE quand un
    utilisateur en anglais télécharge la version française ;
  · le formatage des nombres suit la langue du document. Un séparateur de
    milliers en espace et une virgule décimale sont français ; l'anglais
    attend l'inverse. Le rendre solidaire de la langue évite un document
    anglais parsemé de « 1 234,5 ».
"""

import sys
from pathlib import Path

RACINE = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(RACINE))

# Ce script tourne HORS de `app.py` — en ligne de commande, ou appelé par
# `faire_zip.py`. Il doit donc déclarer lui-même où vivent les traductions :
# le socle refuse de deviner, et le dit clairement plutôt que de servir des
# clés brutes dans un livrable.
from socle import i18n                                          # noqa: E402

i18n.configurer(RACINE / "i18n" / "locales")

from pptx.chart.data import CategoryChartData                   # noqa: E402
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from socle.i18n import LANGUES
from utils.loader import load_all_data
from utils.clean import (
    clean_formations, clean_superieur, clean_indicateur, clean_budget,
)
from socle.stats import econometrie as eco
from utils import analytics


# ─── Charte — celle du SOCLE, dérivée des tokens de l'écran ────────────────
#
# Ces alias gardent les noms courts qu'emploient les dix pages. La charte
# elle-même vit dans `socle.rapport.charte` et CONVERTIT les tokens du tableau
# de bord au lieu de les recopier : la copie précédente laissait le rapport sur
# l'ancienne teinte dès qu'on changeait celle de l'écran, et les deux livrables
# cessaient silencieusement de se ressembler.

from socle.rapport import charte, Langue
from socle.rapport import construire as socle_construire
from socle.rapport import octets as socle_octets
from socle.rapport import generer as socle_generer
from socle.rapport.charte import (
    ENCRE, ENCRE_SECONDAIRE, MUTED, PRIMAIRE, SURFACE, ROUGE,
    SERIE_1, SERIE_2, LARGEUR, HAUTEUR,
)

# Vert du drapeau togolais — couleur de MARQUE, propre à ce défi. Le socle ne
# connaît aucun commanditaire : une teinte de marque se déclare ici.
VERT_TOGO = charte.rgb("#006A4E")

_texte = charte.texte
_titre_page = charte.titre_page
_pied = charte.pied
_bloc_constat = charte.bloc_constat
_bloc_analyse = charte.bloc_analyse
_style_graphe = charte.style_graphe


# ─── Pages ──────────────────────────────────────────────────────────────────

def page_1_couverture(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bande = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.2))
    cadre = bande.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("couv_republique"), taille=14, gras=True,
           couleur=VERT_TOGO, espace_apres=14, premier=True)
    _texte(cadre, lg.t("couv_titre"), taille=40, gras=True, espace_apres=8)
    _texte(cadre, lg.t("couv_sous_titre"),
           taille=17, couleur=ENCRE_SECONDAIRE, espace_apres=26)
    _texte(cadre, lg.t("couv_chiffres", {"total": c["faits"]["total"]}),
           taille=13, couleur=MUTED, espace_apres=6)
    _texte(cadre, lg.t("couv_signature"),
           taille=12, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 1, lg)


def page_2_demarche(prs, c, lg):
    """Données, nettoyage et démarche.

    Trois colonnes plutôt que deux : le nettoyage conditionne la validité de
    tous les résultats qui suivent, il devait figurer dans la présentation et
    non seulement dans le tableau de bord. Le format reste à 10 pages.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "1", lg.t("p1_titre"), lg.t("p1_sous_titre"))

    # ─── Colonne 1 : la question et le matériau ─────────────────────────────
    boite = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(4.0), Inches(5.2))
    cadre = boite.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("p1_question_titre"), taille=14, gras=True,
           espace_apres=5, premier=True)
    _texte(cadre, lg.t("p1_question_corps"),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=14)

    _texte(cadre, lg.t("p1_materiau_titre"), taille=14, gras=True,
           espace_apres=5)
    for index in range(1, 6):
        _texte(cadre, lg.t(f"p1_materiau_{index}"), taille=10,
               couleur=ENCRE_SECONDAIRE, espace_apres=4)

    _texte(cadre, "", taille=5, espace_apres=6)
    _texte(cadre, lg.t("p1_asymetrie_titre"), taille=14, gras=True, espace_apres=5)
    _texte(cadre, lg.t("p1_asymetrie_corps"),
           taille=10, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    # ─── Colonne 2 : le nettoyage ───────────────────────────────────────────
    boite2 = slide.shapes.add_textbox(Inches(4.85), Inches(1.55), Inches(4.0), Inches(5.2))
    cadre2 = boite2.text_frame
    cadre2.word_wrap = True

    _texte(cadre2, lg.t("p1_nettoyage_titre"), taille=14, gras=True,
           espace_apres=3, premier=True)
    _texte(cadre2, lg.t("p1_nettoyage_chapeau"),
           taille=10, couleur=MUTED, espace_apres=8)

    for index in range(1, 7):
        _texte(cadre2, lg.t(f"p1_net{index}_champ"), taille=11, gras=True,
               espace_apres=1)
        _texte(cadre2, lg.t(f"p1_net{index}_traitement"), taille=9.5,
               couleur=ENCRE_SECONDAIRE, espace_apres=7)

    # ─── Colonne 3 : principes et limites ───────────────────────────────────
    boite3 = slide.shapes.add_textbox(Inches(9.1), Inches(1.55), Inches(3.65), Inches(5.2))
    cadre3 = boite3.text_frame
    cadre3.word_wrap = True

    _texte(cadre3, lg.t("p1_principes_titre"), taille=14, gras=True,
           espace_apres=7, premier=True)
    for index in range(1, 4):
        _texte(cadre3, lg.t(f"p1_pr{index}_titre"), taille=11, gras=True,
               espace_apres=1)
        _texte(cadre3, lg.t(f"p1_pr{index}_detail"), taille=9.5,
               couleur=ENCRE_SECONDAIRE, espace_apres=9)

    _texte(cadre3, lg.t("p1_neuvieme_titre"), taille=12,
           gras=True, couleur=PRIMAIRE, espace_apres=3)
    _texte(cadre3, lg.t("p1_neuvieme_corps"),
           taille=9.5, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 2, lg)


def page_3_concentration(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "2", lg.t("p2_titre"), lg.t("p2_sous_titre"))

    faits = c["faits"]
    conc = c["concentration"]

    _bloc_constat(slide, Inches(0.6), Inches(1.55), Inches(2.9), Inches(1.3),
                  lg.pct(faits["part_region_tete"], 1),
                  lg.t("p2_tuile1_libelle", {"region": faits["region_tete"]}),
                  lg.t("p2_tuile1_detail",
                       {"part": lg.nb(faits["part_deux_prefectures"], 1)}),
                  couleur=ROUGE)
    _bloc_constat(slide, Inches(3.7), Inches(1.55), Inches(2.9), Inches(1.3),
                  f"{faits['etab_region_queue']}",
                  lg.t("p2_tuile2_libelle", {"region": faits["region_queue"]}),
                  lg.t("p2_tuile2_detail"), couleur=ROUGE)
    _bloc_constat(slide, Inches(6.8), Inches(1.55), Inches(2.9), Inches(1.3),
                  lg.nb(conc["hhi"], 3), lg.t("p2_tuile3_libelle"),
                  lg.t("p2_tuile3_detail",
                       {"equilibre": lg.nb(conc["hhi_equilibre"], 2)}))
    _bloc_constat(slide, Inches(9.9), Inches(1.55), Inches(2.8), Inches(1.3),
                  lg.nb(conc["gini"], 3), lg.t("p2_tuile4_libelle"),
                  lg.t("p2_tuile4_detail"))

    regions = c["regions"]
    donnees = CategoryChartData()
    donnees.categories = regions["region"].tolist()
    donnees.add_series(lg.t("p2_serie"), regions["etablissements"].tolist())

    graphique = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.6), Inches(3.05),
        Inches(6.1), Inches(3.6), donnees,
    ).chart
    _style_graphe(graphique)
    graphique.has_legend = False
    serie = graphique.plots[0].series[0]
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = SERIE_1
    graphique.plots[0].has_data_labels = True
    graphique.plots[0].data_labels.font.size = Pt(10)

    boite = slide.shapes.add_textbox(Inches(7.0), Inches(3.05), Inches(5.7), Inches(3.6))
    cadre = boite.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("p2_lecture_titre"), taille=14, gras=True,
           espace_apres=6, premier=True)
    _texte(cadre, lg.t("p2_lecture_corps", {
        "tete": faits["region_tete"],
        "part": lg.nb(faits["part_region_tete"], 1),
        "queue": faits["region_queue"],
        "nombre": faits["etab_region_queue"],
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p2_mesure_titre"), taille=14, gras=True, espace_apres=6)
    _texte(cadre, lg.t("p2_mesure_corps", {
        "equilibre": lg.nb(conc["hhi_equilibre"], 2),
        "observe": lg.nb(conc["hhi"], 3),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p2_superieur_titre"), taille=14, gras=True, espace_apres=6)
    _texte(cadre, lg.t("p2_superieur_corps", {
        "lome": lg.nb(c["sup"]["part_lome"], 0),
        "prive": lg.nb(c["sup"]["part_prive"], 0),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 3, lg)


def page_4_dynamique(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "3", lg.t("p3_titre"), lg.t("p3_sous_titre"))

    decennies = c["decennies"]
    donnees = CategoryChartData()
    donnees.categories = decennies["libelle"].tolist()
    donnees.add_series(lg.t("p3_serie"), decennies["creations"].tolist())

    graphique = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.6),
        Inches(7.4), Inches(4.2), donnees,
    ).chart
    _style_graphe(graphique)
    graphique.has_legend = False
    serie = graphique.plots[0].series[0]
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = SERIE_1
    graphique.plots[0].has_data_labels = True
    graphique.plots[0].data_labels.font.size = Pt(10)

    boite = slide.shapes.add_textbox(Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.6))
    cadre = boite.text_frame
    cadre.word_wrap = True

    pic = decennies.loc[decennies["creations"].idxmax()]

    _texte(cadre, lg.t("p3_recente_titre"), taille=14, gras=True, espace_apres=6,
           premier=True)
    _texte(cadre, lg.t("p3_recente_corps", {
        "decennie": int(pic["decennie"]),
        "creations": int(pic["creations"]),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p3_desequilibre_titre"), taille=14, gras=True,
           espace_apres=6)
    _texte(cadre, lg.t("p3_desequilibre_corps"),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p3_inflexion_titre"), taille=14, gras=True,
           couleur=VERT_TOGO, espace_apres=6)
    _texte(cadre, lg.t("p3_inflexion_corps"),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 4, lg)


def page_5_ciseaux(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "4", lg.t("p4_titre"), lg.t("p4_sous_titre"))

    # Les noms de série viennent du domaine `commun` DANS LA LANGUE DU
    # DOCUMENT, pas de ceux que `ciseaux()` a posés : cette fonction traduit
    # pour la session en cours, qui n'est pas forcément celle du fichier
    # demandé — un lecteur en anglais peut télécharger la version française.
    series = c["series_indexees"]
    donnees = CategoryChartData()
    donnees.categories = [int(a) for a in series[0]["x"]]
    donnees.add_series(lg.commun("serie_inscriptions"),
                       [round(v, 1) for v in series[0]["y"]])
    donnees.add_series(lg.commun("serie_depense_etudiant"),
                       [round(v, 1) for v in series[1]["y"]])

    graphique = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.6), Inches(1.6),
        Inches(7.4), Inches(4.2), donnees,
    ).chart
    _style_graphe(graphique)
    graphique.has_legend = True
    graphique.legend.position = XL_LEGEND_POSITION.TOP
    graphique.legend.include_in_layout = False

    for index, couleur in enumerate([SERIE_1, SERIE_2]):
        serie = graphique.plots[0].series[index]
        serie.format.line.color.rgb = couleur
        serie.format.line.width = Pt(2.25)

    boite = slide.shapes.add_textbox(Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.6))
    cadre = boite.text_frame
    cadre.word_wrap = True

    elast = c["elasticite"]

    _texte(cadre, lg.t("p4_trajectoires_titre"), taille=14, gras=True,
           espace_apres=6, premier=True)
    _texte(cadre, lg.t("p4_trajectoires_corps", {
        "facteur": lg.nb(c["acces_indice"] / 100, 1),
        "perte": lg.nb(100 - c["depense_indice"], 0),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p4_axe_titre"), taille=14, gras=True, espace_apres=6)
    _texte(cadre, lg.t("p4_axe_corps"),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p4_ecart_titre"), taille=14, gras=True,
           couleur=PRIMAIRE, espace_apres=6)
    _texte(cadre, lg.t("p4_ecart_corps", {
        "elasticite": lg.nb(abs(elast["elasticite"]), 2),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 5, lg)


def page_6_econometrie(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "5", lg.t("p5_titre"), lg.t("p5_sous_titre"))

    # Calage vertical : 3 blocs à gauche et 2 à droite tiennent entre le filet
    # de titre (1,32") et le bandeau de synthèse (5,15"), sans le chevaucher.
    _bloc_analyse(
        slide, Inches(0.6), Inches(1.52), Inches(6.2), Inches(1.15),
        lg.t("p5_a1_nom"), lg.t("p5_a1_question"),
        lg.t("p5_a1_resultat", {"par_an": lg.nb(c["tendance_ins"]["par_an"], 2)}),
        lg.t("p5_a1_lecture", {"r2": lg.nb(c["tendance_ins"]["r2"] * 100, 0)}),
        couleur_resultat=VERT_TOGO,
    )

    _bloc_analyse(
        slide, Inches(0.6), Inches(2.70), Inches(6.2), Inches(1.15),
        lg.t("p5_a2_nom"), lg.t("p5_a2_question"),
        lg.t("p5_a2_resultat", {"par_an": lg.nb(c["tendance_dep"]["par_an"], 1)}),
        lg.t("p5_a2_lecture"),
        couleur_resultat=ROUGE,
    )

    _bloc_analyse(
        slide, Inches(0.6), Inches(3.88), Inches(6.2), Inches(1.15),
        lg.t("p5_a3_nom"), lg.t("p5_a3_question"),
        lg.t("p5_a3_resultat",
             {"elasticite": lg.nb(c["elasticite"]["elasticite"], 2)}),
        lg.t("p5_a3_lecture",
             {"absolu": lg.nb(abs(c["elasticite"]["elasticite"]), 2)}),
        couleur_resultat=ROUGE,
    )

    # ─── Colonne droite : concentration et rupture ──────────────────────────
    _bloc_analyse(
        slide, Inches(7.1), Inches(1.52), Inches(5.6), Inches(1.75),
        lg.t("p5_a4_nom"), lg.t("p5_a4_question"),
        lg.t("p5_a4_resultat", {"hhi": lg.nb(c["concentration"]["hhi"], 3)}),
        lg.t("p5_a4_lecture", {
            "equilibre": lg.nb(c["concentration"]["hhi_equilibre"], 2),
            "hhi": lg.nb(c["concentration"]["hhi"], 3),
            "gini": lg.nb(c["concentration"]["gini"], 2),
        }),
        couleur_resultat=ROUGE,
    )

    _bloc_analyse(
        slide, Inches(7.1), Inches(3.50), Inches(5.6), Inches(1.55),
        lg.t("p5_a5_nom"), lg.t("p5_a5_question"),
        lg.t("p5_a5_resultat", {"rapport": lg.nb(c["rupture"]["rapport"], 1)}),
        lg.t("p5_a5_lecture", {
            "avant": lg.nb(c["rupture"]["avant"]["pente"], 2),
            "apres": lg.nb(c["rupture"]["apres"]["pente"], 2),
        }),
        couleur_resultat=VERT_TOGO,
    )

    # Bandeau de fiabilité : ce que chaque analyse autorise à conclure.
    etabli = lg.t("p5_verdict_etabli")
    fiabilite = [
        (lg.t("p5_fiab1_nom"), etabli, lg.t("p5_fiab1_detail"), VERT_TOGO),
        (lg.t("p5_fiab2_nom"), etabli, lg.t("p5_fiab2_detail"), VERT_TOGO),
        (lg.t("p5_fiab3_nom"), etabli, lg.t("p5_fiab3_detail"), VERT_TOGO),
        (lg.t("p5_fiab4_nom"), etabli, lg.t("p5_fiab4_detail"), VERT_TOGO),
        (lg.t("p5_fiab5_nom"), lg.t("p5_verdict_non_concluant"),
         lg.t("p5_fiab5_detail", {"n": c["execution"]["n"]}), ROUGE),
    ]

    haut = Inches(5.15)
    for index, (nom, verdict, detail, couleur) in enumerate(fiabilite):
        x = Inches(0.6) + Inches(2.47) * index
        boite = slide.shapes.add_textbox(x, haut, Inches(2.4), Inches(1.2))
        cadre = boite.text_frame
        cadre.word_wrap = True

        _texte(cadre, nom, taille=9, couleur=MUTED, espace_apres=1, premier=True)
        _texte(cadre, verdict, taille=15, gras=True, couleur=couleur, espace_apres=1)
        _texte(cadre, detail, taille=8, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.5))
    cadre = boite.text_frame
    cadre.word_wrap = True
    _texte(cadre, lg.t("p5_note", {"n": c["execution"]["n"]}),
           taille=10, couleur=ENCRE_SECONDAIRE, espace_apres=0, premier=True)

    _pied(slide, 6, lg)


def page_7_budget(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "6", lg.t("p6_titre"), lg.t("p6_sous_titre"))

    compare = c["compare"]
    donnees = CategoryChartData()
    donnees.categories = [int(a) for a in compare["annee"].tolist()]
    donnees.add_series(lg.t("p6_serie_es"), compare["taux_es"].round(1).tolist())
    donnees.add_series(lg.t("p6_serie_educ"), compare["taux_educ"].round(1).tolist())
    donnees.add_series(lg.t("p6_serie_nat"), compare["taux_national"].round(1).tolist())

    graphique = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.6), Inches(1.6),
        Inches(7.4), Inches(4.2), donnees,
    ).chart
    _style_graphe(graphique)
    graphique.has_legend = True
    graphique.legend.position = XL_LEGEND_POSITION.TOP
    graphique.legend.include_in_layout = False

    for index, couleur in enumerate([SERIE_1, SERIE_2, MUTED]):
        serie = graphique.plots[0].series[index]
        serie.format.line.color.rgb = couleur
        serie.format.line.width = Pt(2.25)

    boite = slide.shapes.add_textbox(Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.6))
    cadre = boite.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("p6_lecture_titre"), taille=14, gras=True,
           espace_apres=6, premier=True)
    _texte(cadre, lg.t("p6_lecture_corps", {
        "es": lg.nb(c["moyenne_es"], 1),
        "national": lg.nb(c["moyenne_nat"], 1),
    }), taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p6_alerte_titre"), taille=14, gras=True, espace_apres=6)
    _texte(cadre, lg.t("p6_alerte_corps"),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=12)

    _texte(cadre, lg.t("p6_prive_titre"), taille=14, gras=True, couleur=ROUGE,
           espace_apres=6)
    _texte(cadre, lg.t("p6_prive_corps",
                       {"prive": lg.nb(c["sup"]["part_prive"], 0)}),
           taille=11, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 7, lg)


def page_8_insertion(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "7", lg.t("p7_titre"), lg.t("p7_sous_titre"))

    _bloc_constat(slide, Inches(0.6), Inches(1.6), Inches(2.9), Inches(1.4),
                  lg.pct(c["chomage_fin"]["valeur"], 1),
                  lg.t("p7_tuile1_libelle",
                       {"annee": int(c["chomage_fin"]["annee"])}),
                  lg.t("p7_tuile1_detail"))
    _bloc_constat(slide, Inches(3.7), Inches(1.6), Inches(2.9), Inches(1.4),
                  lg.t("p7_tuile2_valeur"), lg.t("p7_tuile2_libelle"),
                  lg.t("p7_tuile2_detail"), couleur=ROUGE)
    _bloc_constat(slide, Inches(6.8), Inches(1.6), Inches(2.9), Inches(1.4),
                  lg.t("p7_tuile3_valeur"), lg.t("p7_tuile3_libelle"),
                  lg.t("p7_tuile3_detail"), couleur=ROUGE)
    _bloc_constat(slide, Inches(9.9), Inches(1.6), Inches(2.8), Inches(1.4),
                  lg.t("p7_tuile4_valeur"), lg.t("p7_tuile4_libelle"),
                  lg.t("p7_tuile4_detail"), couleur=ROUGE)

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.2))
    cadre = boite.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("p7_paradoxe_titre"), taille=15, gras=True, espace_apres=6,
           premier=True)
    _texte(cadre, lg.t("p7_paradoxe_corps",
                       {"taux": lg.nb(c["chomage_fin"]["valeur"], 1)}),
           taille=12, couleur=ENCRE_SECONDAIRE, espace_apres=14)

    _texte(cadre, lg.t("p7_pilotage_titre"), taille=15, gras=True,
           espace_apres=6)
    _texte(cadre, lg.t("p7_pilotage_corps"),
           taille=12, couleur=ENCRE_SECONDAIRE, espace_apres=14)

    _texte(cadre, lg.t("p7_test_corps"), taille=11, couleur=MUTED, espace_apres=0)

    _pied(slide, 8, lg)


def page_9_leviers(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "8", lg.t("p8_titre"), lg.t("p8_sous_titre"))

    couleurs = [VERT_TOGO, PRIMAIRE, SERIE_2, SERIE_1, ROUGE]

    gauche = Inches(0.6)
    for index, couleur in enumerate(couleurs):
        axe = index + 1
        colonne = index % 3
        rangee = index // 3
        x = gauche + Inches(4.25) * colonne
        y = Inches(1.65) + Inches(2.6) * rangee

        boite = slide.shapes.add_textbox(x, y, Inches(4.0), Inches(2.5))
        cadre = boite.text_frame
        cadre.word_wrap = True

        _texte(cadre, lg.t(f"p8_axe{axe}").upper(), taille=12, gras=True,
               couleur=couleur, espace_apres=5, premier=True)

        for levier in range(1, 5):
            _texte(cadre, f"• {lg.t(f'p8_axe{axe}_l{levier}')}", taille=10,
                   couleur=ENCRE_SECONDAIRE, espace_apres=4)

    _pied(slide, 9, lg)


def page_10_conclusion(prs, c, lg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _titre_page(slide, "9", lg.t("p9_titre"), lg.t("p9_sous_titre"))

    boite = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.1))
    cadre = boite.text_frame
    cadre.word_wrap = True

    _texte(cadre, lg.t("p9_alignement_titre"), taille=16, gras=True,
           espace_apres=10, premier=True)

    for cle, params, couleur in [
        ("p9_territorial", {"region": lg.nb(c["faits"]["part_region_tete"], 0),
                            "lome": lg.nb(c["sup"]["part_lome"], 0)}, ROUGE),
        ("p9_financier", None, ROUGE),
        ("p9_insertion", None, MUTED),
    ]:
        _texte(cadre, lg.t(f"{cle}_titre"), taille=13, gras=True, couleur=couleur,
               espace_apres=3)
        _texte(cadre, lg.t(f"{cle}_detail", params), taille=11,
               couleur=ENCRE_SECONDAIRE, espace_apres=12)

    boite2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.1))
    cadre2 = boite2.text_frame
    cadre2.word_wrap = True

    _texte(cadre2, lg.t("p9_limites_titre"), taille=16, gras=True,
           espace_apres=8, premier=True)
    for index in range(1, 6):
        _texte(cadre2, lg.t(f"p9_limite_{index}"), taille=11,
               couleur=ENCRE_SECONDAIRE, espace_apres=6)

    _texte(cadre2, "", taille=6, espace_apres=8)
    _texte(cadre2, lg.t("p9_decision_titre"), taille=16,
           gras=True, couleur=PRIMAIRE, espace_apres=8)
    _texte(cadre2, lg.t("p9_decision_corps"),
           taille=12, couleur=ENCRE_SECONDAIRE, espace_apres=0)

    _pied(slide, 10, lg)


# ─── Assemblage ─────────────────────────────────────────────────────────────

PAGES = [
    page_1_couverture, page_2_demarche, page_3_concentration,
    page_4_dynamique, page_5_ciseaux, page_6_econometrie,
    page_7_budget, page_8_insertion, page_9_leviers, page_10_conclusion,
]


def collecter():
    """Tous les chiffres de la présentation, calculés depuis les jeux nettoyés."""

    raw = load_all_data()

    formations = clean_formations(raw["formations"])
    superieur = clean_superieur(raw["etablissements"])
    budget = clean_budget(raw["budget"])
    inscriptions = clean_indicateur(raw["inscriptions"])
    depenses = clean_indicateur(raw["depenses"])
    chomage = clean_indicateur(raw["chomage"])

    compare = analytics.execution_comparee(budget)
    series = analytics.ciseaux(inscriptions, depenses, base=1998)
    regions = analytics.par_region(formations)

    return {
        "faits": analytics.concentration(formations),
        "sup": analytics.superieur_synthese(superieur),
        "regions": regions,
        "decennies": analytics.creation_par_decennie(formations),
        "compare": compare,
        "moyenne_es": compare["taux_es"].mean(),
        "moyenne_nat": compare["taux_national"].mean(),
        "series_indexees": series,
        "acces_indice": series[0]["y"][-1],
        "depense_indice": series[1]["y"][-1],
        "chomage_fin": chomage.iloc[-1],
        "tendance_ins": eco.tendance_temporelle(inscriptions),
        "tendance_dep": eco.tendance_temporelle(depenses),
        "elasticite": eco.elasticite(inscriptions, depenses),
        "rupture": eco.rupture_de_tendance(inscriptions, 2000),
        "execution": eco.execution_vs_montant(budget),
        "concentration": eco.concentration(regions["etablissements"].tolist()),
    }


def construire(langue="fr", chiffres=None):
    """La présentation en mémoire, dans la langue demandée.

    Délègue le montage à `socle.rapport` en lui passant NOS dix pages. La
    signature ne change pas : `views/synthese.py` et `scripts/faire_zip.py`
    l'appellent comme avant.
    """

    return socle_construire(PAGES, langue, _chiffres(chiffres))


def _chiffres(fournis):
    """Les chiffres reçus, ou une collecte à la demande.

    Le bouton de téléchargement de la Vue d'ensemble appelle `octets(langue)`
    sans rien passer : il n'a qu'une langue à servir et aucun contexte à
    partager. Le socle, lui, REFUSE de collecter — les agrégations sont du
    code de défi. C'est donc ici que le repli se fait.

    En ligne de commande et dans `faire_zip.py`, les chiffres sont au
    contraire collectés une seule fois puis passés aux deux langues : deux
    collectes laisseraient les documents diverger sur un arrondi.
    """

    return fournis if fournis is not None else collecter()


def octets(langue="fr", chiffres=None):
    """La présentation sérialisée — ce que consomme le bouton de téléchargement."""

    return socle_octets(PAGES, langue, _chiffres(chiffres))


def generer(langue="fr", destination=None, chiffres=None):
    chemin = destination

    if chemin is None:
        lg = Langue(langue)
        chemin = RACINE / "rapport" / f"{lg.t('fichier')}.pptx"

    return socle_generer(PAGES, langue, chemin, _chiffres(chiffres))


if __name__ == "__main__":
    demandees = sys.argv[1:] or list(LANGUES)

    # Les chiffres sont calculés UNE FOIS et partagés : les deux documents
    # doivent porter exactement les mêmes valeurs, et un second chargement des
    # jeux ne servirait qu'à ouvrir la porte à un écart entre les versions.
    partages = collecter()

    for code in demandees:
        chemin, pages = generer(code, chiffres=partages)
        print(f"[{code}] {pages} pages → {chemin}")
